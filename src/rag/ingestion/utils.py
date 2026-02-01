import os
from dataclasses import dataclass
from typing import List, Any

# Use lightweight extractors instead of unstructured (which requires ~4GB of ML dependencies)
from pdfminer.high_level import extract_text as pdf_extract_text
from pdfminer.layout import LAParams
from docx import Document as DocxDocument
from pptx import Presentation
import html

from src.services.llm import openAI
from langchain_core.messages import HumanMessage


@dataclass
class SimpleElement:
    """Simple element to mimic unstructured's element structure"""
    text: str
    element_type: str = "NarrativeText"
    metadata: Any = None

    def __post_init__(self):
        if self.metadata is None:
            self.metadata = SimpleMetadata()


@dataclass  
class SimpleMetadata:
    """Simple metadata to mimic unstructured's metadata structure"""
    page_number: int = 1
    orig_elements: List[Any] = None
    
    def __post_init__(self):
        if self.orig_elements is None:
            self.orig_elements = []


def _extract_pdf(filename: str) -> List[SimpleElement]:
    """Extract text from PDF using pdfminer (lightweight, no ML required)"""
    laparams = LAParams()
    text = pdf_extract_text(filename, laparams=laparams)
    
    # Split by pages (rough approximation - pdfminer doesn't give page breaks easily)
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    
    elements = []
    for i, para in enumerate(paragraphs):
        metadata = SimpleMetadata(page_number=(i // 10) + 1)  # Rough page estimation
        elements.append(SimpleElement(text=para, element_type="NarrativeText", metadata=metadata))
    
    return elements


def _extract_docx(filename: str) -> List[SimpleElement]:
    """Extract text from DOCX using python-docx"""
    doc = DocxDocument(filename)
    elements = []
    
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            element_type = "Title" if para.style and "heading" in para.style.name.lower() else "NarrativeText"
            metadata = SimpleMetadata(page_number=1)
            elements.append(SimpleElement(text=para.text.strip(), element_type=element_type, metadata=metadata))
    
    # Also extract tables
    for table in doc.tables:
        table_html = "<table>"
        for row in table.rows:
            table_html += "<tr>"
            for cell in row.cells:
                table_html += f"<td>{html.escape(cell.text)}</td>"
            table_html += "</tr>"
        table_html += "</table>"
        
        table_element = SimpleElement(text=table_html, element_type="Table", metadata=SimpleMetadata(page_number=1))
        elements.append(table_element)
    
    return elements


def _extract_pptx(filename: str) -> List[SimpleElement]:
    """Extract text from PPTX using python-pptx"""
    prs = Presentation(filename)
    elements = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                metadata = SimpleMetadata(page_number=slide_num)
                elements.append(SimpleElement(text=shape.text.strip(), element_type="NarrativeText", metadata=metadata))
            
            # Handle tables in slides
            if shape.has_table:
                table = shape.table
                table_html = "<table>"
                for row in table.rows:
                    table_html += "<tr>"
                    for cell in row.cells:
                        table_html += f"<td>{html.escape(cell.text)}</td>"
                    table_html += "</tr>"
                table_html += "</table>"
                
                table_element = SimpleElement(text=table_html, element_type="Table", metadata=SimpleMetadata(page_number=slide_num))
                elements.append(table_element)
    
    return elements


def _extract_text(filename: str) -> List[SimpleElement]:
    """Extract from plain text file"""
    with open(filename, 'r', encoding='utf-8', errors='ignore') as f:
        text = f.read()
    
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    return [SimpleElement(text=para, element_type="NarrativeText") for para in paragraphs]


def _extract_md(filename: str) -> List[SimpleElement]:
    """Extract from markdown file"""
    return _extract_text(filename)  # Treat as plain text for now


def _extract_html(filename: str) -> List[SimpleElement]:
    """Extract from HTML file"""
    with open(filename, 'r', encoding='utf-8', errors='ignore') as f:
        content = f.read()
    
    # Simple HTML text extraction (strip tags)
    import re
    text = re.sub(r'<[^>]+>', ' ', content)
    text = re.sub(r'\s+', ' ', text).strip()
    
    return [SimpleElement(text=text, element_type="NarrativeText")]


def partition_document(temp_file: str, file_type: str, source_type: str = "file"):
    """Partition document based on file type and source type using lightweight extractors"""
    source = (source_type or "file").lower()
    
    if source == "url":
        return _extract_html(temp_file)

    kind = (file_type or "").lower()
    dispatch = {
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "pptx": _extract_pptx,
        "txt": _extract_text,
        "md": _extract_md,
    }

    if kind not in dispatch:
        raise ValueError(f"Unsupported file_type: {file_type}")

    return dispatch[kind](temp_file)    


@dataclass
class SimpleChunk:
    """Represents a chunk of content"""
    text: str
    metadata: Any = None
    
    def __post_init__(self):
        if self.metadata is None:
            self.metadata = SimpleChunkMetadata()


@dataclass
class SimpleChunkMetadata:
    """Metadata for a chunk"""
    page_number: int = 1
    orig_elements: List[Any] = None
    
    def __post_init__(self):
        if self.orig_elements is None:
            self.orig_elements = []


def chunk_by_title(
    elements: List[SimpleElement],
    max_characters: int = 3000,
    new_after_n_chars: int = 2400,
    combine_text_under_n_chars: int = 500,
) -> List[SimpleChunk]:
    """
    Lightweight implementation of chunk_by_title that mimics unstructured's behavior.
    Groups elements together respecting character limits.
    """
    if not elements:
        return []
    
    chunks = []
    current_texts = []
    current_elements = []
    current_char_count = 0
    current_page = 1
    
    for element in elements:
        element_text = element.text if hasattr(element, 'text') else str(element)
        element_len = len(element_text)
        
        # Get page number from element if available
        if hasattr(element, 'metadata') and hasattr(element.metadata, 'page_number'):
            current_page = element.metadata.page_number
        
        # Check if this element is a title/header (signals new section)
        is_title = False
        if hasattr(element, 'element_type'):
            is_title = element.element_type in ["Title", "Header"]
        
        # Decide if we should start a new chunk
        should_new_chunk = False
        
        # Start new chunk if this is a title and current chunk is substantial
        if is_title and current_char_count >= combine_text_under_n_chars:
            should_new_chunk = True
        
        # Start new chunk if we've exceeded the soft limit
        if current_char_count >= new_after_n_chars:
            should_new_chunk = True
        
        # Start new chunk if adding this would exceed hard limit
        if current_char_count + element_len > max_characters and current_texts:
            should_new_chunk = True
        
        # Create new chunk if needed
        if should_new_chunk and current_texts:
            chunk_text = "\n\n".join(current_texts)
            chunk_metadata = SimpleChunkMetadata(
                page_number=current_page,
                orig_elements=current_elements.copy()
            )
            chunks.append(SimpleChunk(text=chunk_text, metadata=chunk_metadata))
            current_texts = []
            current_elements = []
            current_char_count = 0
        
        # Add element to current chunk
        current_texts.append(element_text)
        current_elements.append(element)
        current_char_count += element_len
    
    # Don't forget the last chunk
    if current_texts:
        chunk_text = "\n\n".join(current_texts)
        chunk_metadata = SimpleChunkMetadata(
            page_number=current_page,
            orig_elements=current_elements.copy()
        )
        chunks.append(SimpleChunk(text=chunk_text, metadata=chunk_metadata))
    
    # Combine small chunks if they're under the threshold
    if len(chunks) > 1:
        merged_chunks = []
        i = 0
        while i < len(chunks):
            current = chunks[i]
            
            # If current chunk is small, try to merge with next
            while (i + 1 < len(chunks) and 
                   len(current.text) < combine_text_under_n_chars and
                   len(current.text) + len(chunks[i + 1].text) <= max_characters):
                next_chunk = chunks[i + 1]
                current = SimpleChunk(
                    text=current.text + "\n\n" + next_chunk.text,
                    metadata=SimpleChunkMetadata(
                        page_number=current.metadata.page_number,
                        orig_elements=current.metadata.orig_elements + next_chunk.metadata.orig_elements
                    )
                )
                i += 1
            
            merged_chunks.append(current)
            i += 1
        
        chunks = merged_chunks
    
    return chunks    


def analyze_elements(elements):
    """Analyze the elements and return the summary"""

    text_count = 0
    table_count = 0
    image_count = 0
    title_count = 0
    other_count = 0

    # Go through each element and count what type it is
    for element in elements:
        element_name = type(
            element
        ).__name__  # __name__ is a special attribute that returns the class name like "Table" or "NarrativeText"

        if element_name == "Table":
            table_count += 1
        elif element_name == "Image":
            image_count += 1
        elif element_name in ["Title", "Header"]:
            title_count += 1
        elif element_name in ["NarrativeText", "Text", "ListItem", "FigureCaption"]:
            text_count += 1
        else:
            other_count += 1

    # Return a simple dictionary
    return {
        "text": text_count,
        "tables": table_count,
        "images": image_count,
        "titles": title_count,
        "other": other_count,
    }


def separate_content_types(chunk, source_type="file"):
    """Analyze what types of content are in a chunk"""
    is_url_source = source_type == "url"

    content_data = {
        "text": chunk.text,  # By default every chunk will have text so chunk.text will not be None.
        "tables": [],
        "images": [],
        "types": ["text"],
    }

    # Check for tables and images in original elements
    if hasattr(chunk, "metadata") and hasattr(chunk.metadata, "orig_elements"):
        # orig_elements list all the atomic elements in the chunk.
        for element in chunk.metadata.orig_elements:
            element_type = type(element).__name__

            # Handle tables
            if element_type == "Table":
                content_data["types"].append("table")
                # getattr is a built-in function that returns the value of the named attribute of an object.
                # text_as_html will return the HTML representation of the table if it exists, otherwise it will return the text attribute of the element.
                table_html = getattr(element.metadata, "text_as_html", element.text)
                content_data["tables"].append(table_html)

            # Handle images (skip for URL sources)
            elif element_type == "Image" and not is_url_source:
                if (
                    hasattr(element, "metadata")
                    and hasattr(element.metadata, "image_base64")
                    and element.metadata.image_base64 is not None
                ):
                    content_data["types"].append("image")
                    content_data["images"].append(element.metadata.image_base64)

    content_data["types"] = list(set(content_data["types"]))

    # Return the content data (Below is the example return structure)
    # {
    #     "text": "This is the main text content of the chunk...",
    #     "tables": ["<table><tr><th>Header</th></tr><tr><td>Data</td></tr></table>"],
    #     "images": ["iVBORw0KGgoAAAANSUhEUgAA..."],  # base64 encoded image strings
    #     "types": ["text", "table", "image"]  # or ["text"], ["text", "table"], etc.
    # }

    return content_data

def get_page_number(chunk, chunk_index):
    """Get page number from chunk or use fallback"""
    if hasattr(chunk, "metadata"):
        page_number = getattr(chunk.metadata, "page_number", None)
        if page_number is not None:
            return page_number


    return chunk_index + 1



def create_ai_summary(text, tables_html, images_base64):
    """Create AI-enhanced summary for tables and images present in the chunks"""

    try:
        # Build the text prompt with more efficient instructions
        prompt_text = f"""
            Create a searchable index for this document content.
            CONTENT:
            {text}
        """

        # Add tables if present
        if tables_html:
            prompt_text += "TABLES:\n"
            for i, table in enumerate(tables_html):
                prompt_text += f"Table {i+1}:\n{table}\n\n"

        # More concise but effective prompt
        prompt_text += """
            Generate a structured search index (aim for 250-400 words):

            QUESTIONS: List 5-7 key questions this content answers (use what/how/why/when/who variations)

            KEYWORDS: Include:
            - Specific data (numbers, dates, percentages, amounts)
            - Core concepts and themes
            - Technical terms and casual alternatives
            - Industry terminology

            VISUALS (if images present):
            - Chart/graph types and what they show
            - Trends and patterns visible
            - Key insights from visualizations

            DATA RELATIONSHIPS (if tables present):
            - Column headers and their meaning
            - Key metrics and relationships
            - Notable values or patterns

            Focus on terms users would actually search for. Be specific and comprehensive.

            SEARCH INDEX:"""

        # Build message content starting with the text prompt
        message_content = [{"type": "text", "text": prompt_text}]

        # Add images to the message
        for i, image_base64 in enumerate(images_base64):
            message_content.append(
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{image_base64}"},
                }
            )
            # print(f"🖼️ Image {i+1} included in summary request")

        message = HumanMessage(content=message_content)
        response = openAI["embeddings_llm"].invoke([message])

        return response.content

    except Exception as e:
        raise Exception(f"Failed to create AI summary: {str(e)}")